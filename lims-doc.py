import os
import PIL.Image as PILImage
import PIL.ImageDraw as PILDraw
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

os.makedirs("generated_diagrams", exist_ok=True)

# 1. Struktur Data Diagram Alur Visual (Kanan Slide)
steps_diagram_data = [
    (1, "DIAGRAM PROSES STEP 1", [
        ("Input Sampel Baru", "Form LIMS diisi metadata fisik sampel & pengirim"),
        ("Generate Barcode ID", "Sistem menerbitkan ID unik & cetak label fisik"),
        ("Tempel & Simpan Temp", "Barcode ditempel di sampel, status 'Registered'")
    ]),
    (2, "DIAGRAM PROSES STEP 2", [
        ("Mapping Parameter Uji", "Penentuan item pengujian & metode baku"),
        ("Splitting Alikot", "Pembagian porsi sampel untuk tiap departemen uji"),
        ("Assign Worklist Analis", "Distribusi tugas ke antrean kerja laboratorium")
    ]),
    (3, "DIAGRAM PROSES STEP 3", [
        ("Scan Barcode Sampel", "Konfirmasi penerimaan sampel di meja preparasi"),
        ("Ekstraksi & Reagensia", "Pelaksanaan preparasi & logging lot reagen"),
        ("Ready for Analysis", "Sampel siap diukur, status 'Prepared'")
    ]),
    (4, "DIAGRAM PROSES STEP 4", [
        ("Grouping Run Batch", "Pengelompokan sampel dalam satu urutan uji"),
        ("Insert QC Controls", "Penambahan sampel kontrol untuk penjaminan mutu"),
        ("Finalizing Sequence", "Matriks urutan injeksi siap dikirim ke alat")
    ]),
    (5, "DIAGRAM PROSES STEP 5", [
        ("Run Instrument", "Alat menjalankan pengujian sampel otomatis"),
        ("Capture Raw Data", "LIMS menyedot file hasil/kromatogram secara direct"),
        ("Environmental Check", "Log otomatis kondisi ruang lab saat analisis")
    ]),
    (6, "DIAGRAM PROSES STEP 6", [
        ("Apply Calculation Formula", "Pengolahan nilai mentah dengan rumus standar"),
        ("Evaluate QC Limits", "Pengecekan otomatis batas presisi & akurasi QC"),
        ("Auto-Flagging Status", "Penandaan status Normal vs OOS/OOT")
    ]),
    (7, "DIAGRAM PROSES STEP 7", [
        ("Supervisor Data Check", "Peninjauan data mentah & grafik kalibrasi"),
        ("OOS Handling (If Any)", "Keputusan re-sample, re-test, atau lanjut"),
        ("Technical Sign-Off", "Otentikasi tingkat pertama oleh teknisi senior")
    ]),
    (8, "DIAGRAM PROSES STEP 8", [
        ("Compile Certificate/LHA", "Penggabungan data parameter ke layout sertifikat"),
        ("Digital Signature", "Pengesahan dokumen oleh Kepala Laboratorium"),
        ("Final Data Locking", "Penguncian data agar bersifat permanent record")
    ]),
    (9, "DIAGRAM PROSES STEP 9", [
        ("Send LHA to Client", "Publikasi laporan ke portal pelanggan / e-mail"),
        ("Sample Storage / Disposal", "Penyimpanan sisa sampel sesuai masa retensi"),
        ("Audit Trail Archiving", "Penyimpanan log aktivitas lengkap untuk audit ISO")
    ])
]

# Generate Gambar Diagram
diagram_file_paths = []
for step_num, header_title, cards in steps_diagram_data:
    width, height = 825, 855
    img = PILImage.new('RGB', (width, height), color=(248, 250, 252))
    draw = PILDraw.Draw(img)
    
    color_teal = (14, 116, 144)
    color_border = (226, 232, 240)
    color_card_bg = (255, 255, 255)
    color_text_muted = (80, 95, 110)
    
    draw.rounded_rectangle([15, 15, width-15, height-15], radius=35, fill=(248, 250, 252), outline=color_teal, width=5)
    draw.rounded_rectangle([45, 45, width-45, 135], radius=15, fill=color_teal)
    draw.text((215, 72), header_title, fill=(255, 255, 255))
    
    for idx, (node_title, node_desc) in enumerate(cards):
        top_y = 175 + (idx * 215)
        bot_y = top_y + 180
        draw.rounded_rectangle([55, top_y, width-55, bot_y], radius=20, fill=color_card_bg, outline=color_border, width=3)
        draw.polygon([(85, top_y + 40), (85, top_y + 60), (105, top_y + 50)], fill=color_teal)
        draw.text((120, top_y + 35), node_title, fill=color_teal)
        draw.text((85, top_y + 90), node_desc, fill=color_text_muted)
        
    file_p = f"generated_diagrams/step_{step_num}_diagram.png"
    img.save(file_p)
    diagram_file_paths.append(file_p)

# 2. Inisialisasi PowerPoint Presentasi
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

COLOR_BG = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(20, 30, 45)
COLOR_TEXT_MUTED = RGBColor(80, 95, 110)
COLOR_PRIMARY = RGBColor(14, 116, 144)
COLOR_CARD_BG = RGBColor(248, 250, 252)
COLOR_CARD_BORDER = RGBColor(226, 232, 240)

def set_slide_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLOR_BG

def create_diagram_slide(prs, step_num, title, subtitle, bullets, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    
    badge_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(6.0), Inches(0.4))
    p_badge = badge_box.text_frame.paragraphs[0]
    p_badge.text = f"PROSES STEP {step_num} DARI 9"
    p_badge.font.size = Pt(12)
    p_badge.font.bold = True
    p_badge.font.color.rgb = COLOR_PRIMARY
    
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(11.73), Inches(1.1))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(26)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_DARK
    
    p_sub = tf_title.add_paragraph()
    p_sub.text = subtitle
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = COLOR_PRIMARY
    
    card_left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(5.8), Inches(4.5))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = COLOR_CARD_BG
    card_left.line.color.rgb = COLOR_CARD_BORDER
    
    content_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.3), Inches(5.4), Inches(4.3))
    tf_content = content_box.text_frame
    tf_content.word_wrap = True
    
    for idx, (b_title, b_desc) in enumerate(bullets):
        p = tf_content.paragraphs[0] if idx == 0 else tf_content.add_paragraph()
        p.space_after = Pt(12)
        run_t = p.add_run()
        run_t.text = f"{idx+1}. {b_title}\n"
        run_t.font.bold = True
        run_t.font.size = Pt(13)
        run_t.font.color.rgb = COLOR_TEXT_DARK
        run_d = p.add_run()
        run_d.text = f"   {b_desc}"
        run_d.font.size = Pt(11)
        run_d.font.color.rgb = COLOR_TEXT_MUTED

    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(7.0), Inches(2.2), width=Inches(5.5), height=Inches(4.5))

process_slides_data = [
    (1, "Register LIMS & Barcode", "Tahap 1: Pengumpulan & Perencanaan Sampel", [("Pencatatan Sampel Masuk", "Registrasi sampel ke sistem LIMS lengkap dengan informasi metadata penerimaan."), ("Penjanaan Barcode Unik", "Sistem mencetak label barcode unik untuk penempelan pada wadah sampel."), ("Pencatatan Identitas Pengirim", "Identitas pengirim, tanggal terima, dan kondisi fisik sampel tercatat otomatis.")]),
    (2, "Alokasi & Distribusi Sampel", "Tahap 1: Pengumpulan & Perencanaan Sampel", [("Penentuan Parameter Uji", "Sistem memetakan sampel ke jenis pengujian yang diperlukan sesuai permintaan."), ("Pembagian Alikot Sampel", "Sampel dibagi ke dalam beberapa wadah kecil (alikot) jika membutuhkan beberapa lab/alat."), ("Penjadwalan Kerja (Worklist)", "Sampel secara otomatis masuk ke antrean kerja analis atau instrumen terkait.")]),
    (3, "Preparation & Extraction", "Tahap 1: Pengumpulan & Perencanaan Sampel", [("Verifikasi Sampel Pra-Uji", "Analis melakukan pemindaian barcode sampel untuk mengonfirmasi kesiapan preparasi."), ("Proses Penyiapan Kimia/Fisik", "Ekstraksi, destruksi, atau pelarutan sampel sesuai standar operasional prosedur (SOP)."), ("Pencatatan Reagen & Matriks", "Input otomatis konsumsi bahan kimia, lot number, dan tanggal kadaluarsa reagen.")]),
    (4, "Batching & QC Blank/Spike", "Tahap 2: Eksekusi Pengujian Laboratorium", [("Pengelompokan Run (Batching)", "Pengelompokan beberapa sampel ke dalam satu run pengujian agar efisien."), ("Inserksi QC Control", "Penambahan sampel kontrol untuk penjaminan mutu."), ("Validasi Antrean Uji", "Sistem memvalidasi urutan posisi sampel sebelum dimasukkan ke dalam alat analit.")]),
    (5, "Instrument Interfacing & Data Capture", "Tahap 2: Eksekusi Pengujian Laboratorium", [("Direct Instrument Interfacing", "Penyerapan data hasil otomatis dari alat analit (HPLC, GC, Spektro) ke LIMS."), ("Logging Kondisi Lingkungan", "Pencatatan suhu dan kelembaban ruang lab secara real-time saat pengujian."), ("Reduksi Risiko Human Error", "Mengeliminasi ketik manual angka hasil dari layar alat ke lembar kerja.")]),
    (6, "Kalkulasi Hasil & Auto-Validation", "Tahap 2: Eksekusi Pengujian Laboratorium", [("Kalkulasi Formula Otomatis", "Perhitungan konsentrasi, faktor pengenceran, dan konversi unit berbasis rumus LIMS."), ("Evaluasi QC Curve", "Pembacaan kurva kalibrasi dan batas keberhasilan QC sampel secara presisi."), ("Flagging Limit Spesifikasi", "Sistem menandai secara otomatis hasil yang Out of Specification (OOS).")]),
    (7, "Technical Review & Peer Approval", "Tahap 3: Validasi & Pelaporan Hasil", [("Verifikasi oleh Supervisor", "Pemeriksaan kembali kurva, nilai r-square, dan keabsahan prosedur oleh Senior Uji."), ("Penanganan Out of Spec (OOS)", "Pemicuan alur kerja re-test atau investigasi jika ditemukan ketidaksesuaian data."), ("Otorisasi Tingkat Pertama", "Persetujuan teknis (technical sign-off) sebelum data diteruskan ke penyusun LHA.")]),
    (8, "Penerbitan LHA / COA", "Tahap 3: Validasi & Pelaporan Hasil", [("Kompilasi Sertifikat Hasil", "Penggabungan seluruh parameter uji ke dalam template Laporan Hasil Analisis (LHA)."), ("Digital Signature Manager", "Penandatanganan dokumen secara digital berenkripsi oleh Manajer Mutu/Lab."), ("Penguncian Data (Data Locking)", "Pembekuan rekam medis uji agar tidak dapat diubah tanpa izin revisi khusus.")]),
    (9, "Archiving & Sample Retention", "Tahap 3: Validasi & Pelaporan Hasil", [("Pengiriman LHA ke Klien", "Dokumen dikirim otomatis via portal pelanggan, e-mail, atau integrasi sistem API."), ("Manajemen Retensi Sampel", "Penyimpanan sisa sampel di cold storage dengan jadwal pemusnahan otomatis."), ("Arsip Audit Trail", "Penyimpanan seluruh jejak aktivitas pengguna dan log data untuk kebutuhan akreditasi.")])
]

for idx, (step_num, title, subtitle, bullets) in enumerate(process_slides_data):
    create_diagram_slide(prs, step_num, title, subtitle, bullets, diagram_file_paths[idx])

prs.save("Presentasi_LIMS_9_Step_Visual.pptx")
print("File PPTX Presentasi_LIMS_9_Step_Visual.pptx berhasil dibuat!")